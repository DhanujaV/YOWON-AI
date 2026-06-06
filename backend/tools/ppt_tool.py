"""
tools/ppt_tool.py — PowerPoint presentation extractor using python-pptx.

Extracts per-slide:
  - Slide title
  - All text content
  - Notes (speaker notes)
  - Shape types (for detecting diagrams)
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_ppt_data(ppt_path: str | Path) -> dict[str, Any]:
    """
    Parse a .pptx file and return structured slide content.

    Returns:
      {
        "slide_count": int,
        "slides": list of {
            "slide_number": int,
            "title": str,
            "text": str,       # all text on the slide
            "notes": str,      # speaker notes
            "has_diagram": bool
        },
        "full_text": str
      }
    """
    path = Path(ppt_path)
    if not path.exists():
        return {"error": f"File not found: {path}"}

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        return {"error": f"Could not open PPT: {exc}"}

    slides_data: list[dict] = []
    all_text_parts: list[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        title_text = ""
        slide_text_parts: list[str] = []
        has_diagram = False

        for shape in slide.shapes:
            # ── Detect diagrams / SmartArt / charts ───────────────────────
            if shape.shape_type in (
                MSO_SHAPE_TYPE.GROUP,        # 6
                MSO_SHAPE_TYPE.PICTURE,      # 13
            ):
                has_diagram = True

            # Check for chart or SmartArt via shape name heuristic
            name_lower = shape.name.lower()
            if any(kw in name_lower for kw in ("chart", "diagram", "smartart", "process")):
                has_diagram = True

            # ── Extract text ──────────────────────────────────────────────
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                line = " ".join(run.text for run in paragraph.runs).strip()
                if not line:
                    continue
                slide_text_parts.append(line)

                # First text of the title placeholder is the slide title
                if shape.name.lower().startswith("title") and not title_text:
                    title_text = line

        # ── Speaker notes ─────────────────────────────────────────────────
        notes_text = ""
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_text = notes_frame.text.strip() if notes_frame else ""

        slide_combined = "\n".join(slide_text_parts)
        slides_data.append({
            "slide_number": slide_num,
            "title": title_text or f"Slide {slide_num}",
            "text": slide_combined,
            "notes": notes_text,
            "has_diagram": has_diagram,
        })
        all_text_parts.append(f"=== Slide {slide_num}: {title_text} ===\n{slide_combined}")

    return {
        "slide_count": len(slides_data),
        "slides": slides_data,
        "full_text": "\n\n".join(all_text_parts),
    }